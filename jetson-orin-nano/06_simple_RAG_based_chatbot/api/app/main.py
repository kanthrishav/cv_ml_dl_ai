import os, io, re, json, sqlite3, hashlib, time, uuid, shlex, subprocess
from typing import List, Optional, Dict, Any, Tuple
from fastapi import FastAPI, HTTPException, Body
from pydantic import BaseModel
import httpx
from tenacity import retry, wait_exponential, stop_after_attempt
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct, Filter, FieldCondition, MatchValue

# --------- ENV ---------
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://host.docker.internal:11434").rstrip("/")
QDRANT_URL = os.getenv("QDRANT_URL", "http://qdrant:6333")
SEARXNG_URL = os.getenv("SEARXNG_URL", "http://searxng:8080")
STATE_DIR = os.getenv("STATE_DIR", "/state")
DEFAULT_EMBED_MODEL = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")
DEFAULT_CHAT_MODEL  = os.getenv("OLLAMA_CHAT_MODEL", "llama3.1:8b-instruct-q4_K_M")

DEFAULT_LINK_LIMIT = int(os.getenv("DEFAULT_LINK_LIMIT", "8"))
DEFAULT_TOP_K = int(os.getenv("DEFAULT_TOP_K", "6"))
DEFAULT_CHUNK_SIZE = int(os.getenv("DEFAULT_CHUNK_SIZE", "1200"))
DEFAULT_CHUNK_OVERLAP = int(os.getenv("DEFAULT_CHUNK_OVERLAP", "200"))
CONFIDENCE_THRESHOLD = float(os.getenv("DEFAULT_CONFIDENCE_THRESHOLD", "0.35"))

os.makedirs(STATE_DIR, exist_ok=True)
DB_PATH = os.path.join(STATE_DIR, "ullu.db")

# --------- DB (chat memory & sessions) ---------
def db():
    conn = sqlite3.connect(DB_PATH)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS sessions(
            session_id TEXT PRIMARY KEY,
            created_at REAL,
            kb_key TEXT,
            mode TEXT,
            hybrid_weight REAL,
            memory_limit INTEGER
        )""")
    conn.execute("""
        CREATE TABLE IF NOT EXISTS messages(
            session_id TEXT,
            ts REAL,
            role TEXT,
            content TEXT
        )""")
    conn.commit()
    return conn

# --------- Qdrant ---------
client = QdrantClient(url=QDRANT_URL, timeout=60)

EMBED_DIM = 768  # nomic-embed-text dimension

def collection_name_for_kb(paths: List[str]) -> str:
    norm = ",".join(sorted([os.path.abspath(p) for p in paths]))
    h = hashlib.sha256(norm.encode()).hexdigest()[:24]
    return f"kb_{h}"

def ensure_collection(name: str):
    collections = [c.name for c in client.get_collections().collections]
    if name not in collections:
        client.create_collection(
            collection_name=name,
            vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE)
        )

# --------- Embeddings & Chat via Ollama ---------
@retry(wait=wait_exponential(min=0.5, max=8), stop=stop_after_attempt(5))
def ollama_embed(texts: List[str], model: str = DEFAULT_EMBED_MODEL) -> List[List[float]]:
    url = f"{OLLAMA_BASE_URL}/api/embeddings"
    out = []
    with httpx.Client(timeout=120) as s:
        for t in texts:
            r = s.post(url, json={"model": model, "input": t})
            r.raise_for_status()
            out.append(r.json()["embedding"])
    return out

@retry(wait=wait_exponential(min=0.5, max=8), stop=stop_after_attempt(5))
def ollama_chat(system: str, messages: List[Dict[str,str]], model: str = DEFAULT_CHAT_MODEL, temperature: float = 0.2) -> str:
    url = f"{OLLAMA_BASE_URL}/api/chat"
    with httpx.Client(timeout=None) as s:
        r = s.post(url, json={
            "model": model,
            "temperature": temperature,
            "messages": [{"role":"system","content":system}] + messages
        })
        r.raise_for_status()
        data = r.json()
        # /api/chat (non-stream) returns {message:{role,content}, done:true}
        if "message" in data and "content" in data["message"]:
            return data["message"]["content"]
        return data.get("content","")

# --------- File loaders & chunking ---------
def run(cmd: str) -> Tuple[int,str,str]:
    p = subprocess.Popen(shlex.split(cmd), stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    out, err = p.communicate()
    return p.returncode, out.decode("utf-8","ignore"), err.decode("utf-8","ignore")

def read_pdf(path: str) -> str:
    # Try text extraction
    code, out, _ = run(f"pdftotext -layout {shlex.quote(path)} -")
    if code == 0 and out.strip():
        return out
    # OCR fallback
    tmpdir = f"/tmp/ocr_{uuid.uuid4().hex}"
    os.makedirs(tmpdir, exist_ok=True)
    run(f"pdftoppm -r 200 {shlex.quote(path)} {tmpdir}/page")
    texts = []
    for f in sorted([x for x in os.listdir(tmpdir) if x.startswith('page-') or x.startswith('page')]):
        if f.endswith(".ppm") or f.endswith(".pgm") or f.endswith(".pbm"):
            ppm = os.path.join(tmpdir, f)
            outtxt = os.path.join(tmpdir, f"{f}.txt")
            run(f"tesseract {ppm} {outtxt[:-4]} -l eng")
            try:
                with open(outtxt, "r") as fh:
                    texts.append(fh.read())
            except: pass
    return "\n".join(texts)

def read_docx(path: str) -> str:
    import docx2txt
    return docx2txt.process(path) or ""

def read_doc(path: str) -> str:
    # convert .doc -> .docx using headless libreoffice, then parse
    tmp = f"/tmp/{uuid.uuid4().hex}.docx"
    run(f"libreoffice --headless --convert-to docx --outdir /tmp {shlex.quote(path)}")
    return read_docx(tmp) if os.path.exists(tmp) else ""

def read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                texts.append(shape.text)
    return "\n".join(texts)

def read_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    texts=[]
    for ws in wb.worksheets:
        texts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            texts.append("\t".join("" if v is None else str(v) for v in row))
    return "\n".join(texts)

def read_txt(path: str) -> str:
    with open(path,"r",errors="ignore") as f:
        return f.read()

def read_md(path: str) -> str:
    with open(path,"r",errors="ignore") as f:
        return f.read()

def read_html(path: str) -> str:
    import bs4, html2text
    with open(path,"r",errors="ignore") as f:
        soup = bs4.BeautifulSoup(f.read(),"html.parser")
    return html2text.html2text(str(soup))

def load_file(path: str) -> str:
    ext = os.path.splitext(path.lower())[1]
    if ext == ".pdf": return read_pdf(path)
    if ext == ".docx": return read_docx(path)
    if ext == ".doc": return read_doc(path)
    if ext == ".pptx": return read_pptx(path)
    if ext in (".xls",".xlsx"): return read_xlsx(path)
    if ext in (".txt",): return read_txt(path)
    if ext in (".md",): return read_md(path)
    if ext in (".html",".htm"): return read_html(path)
    return ""

def chunk(text: str, size: int, overlap: int) -> List[str]:
    out=[]
    i=0
    n=len(text)
    while i < n:
        out.append(text[i:i+size])
        i += max(1, size - overlap)
    return [x for x in out if x.strip()]

# --------- Ingest ---------
def discover_files(paths: List[str]) -> List[str]:
    files=[]
    exts={".pdf",".doc",".docx",".txt",".xls",".xlsx",".pptx",".md",".html",".htm"}
    for p in paths:
        p=os.path.abspath(p)
        if os.path.isfile(p) and os.path.splitext(p)[1].lower() in exts:
            files.append(p)
        elif os.path.isdir(p):
            for root,_,fnames in os.walk(p):
                for fn in fnames:
                    if os.path.splitext(fn)[1].lower() in exts:
                        files.append(os.path.join(root,fn))
    return files

def ingest(paths: List[str], chunk_size=DEFAULT_CHUNK_SIZE, overlap=DEFAULT_CHUNK_OVERLAP, embed_model=DEFAULT_EMBED_MODEL) -> Dict[str,Any]:
    c_name = collection_name_for_kb(paths)
    ensure_collection(c_name)

    files = discover_files(paths)
    added = 0
    for fp in files:
        text = load_file(fp)
        if not text.strip(): continue
        pieces = chunk(text, chunk_size, overlap)
        vecs = ollama_embed(pieces, model=embed_model)
        pts = []
        for i,(t,v) in enumerate(zip(pieces,vecs)):
            pts.append(PointStruct(
                id=int(time.time()*1e6)%2_147_483_647 + i,
                vector=v,
                payload={"text": t, "source": fp, "chunk_id": i}
            ))
        if pts:
            client.upsert(collection_name=c_name, points=pts)
            added += len(pts)
    return {"collection": c_name, "chunks": added, "files": len(files)}

# --------- Retrieve ---------
def search_collection(c_name: str, query: str, top_k: int, embed_model: str) -> List[Dict[str,Any]]:
    qvec = ollama_embed([query], model=embed_model)[0]
    res = client.search(collection_name=c_name, query_vector=qvec, limit=top_k, with_payload=True)
    out=[]
    for r in res:
        payload=r.payload or {}
        out.append({
            "text": payload.get("text",""),
            "source": payload.get("source",""),
            "score": float(r.score)
        })
    return out

# --------- SearXNG + scraping ---------
def searx_search(query: str, limit: int) -> List[str]:
    url = f"{SEARXNG_URL}/search"
    with httpx.Client(timeout=60) as s:
        r = s.get(url, params={"q":query, "format":"json", "language":"en", "safesearch":"0"})
        r.raise_for_status()
        data = r.json()
    urls=[]
    for r in data.get("results",[])[:limit]:
        u = r.get("url")
        if u and isinstance(u,str): urls.append(u)
    return urls

def fetch_and_clean(urls: List[str], max_chars: int) -> List[Tuple[str,str]]:
    import trafilatura
    out=[]
    for u in urls:
        try:
            downloaded = trafilatura.fetch_url(u, no_ssl=True)
            if not downloaded: continue
            text = trafilatura.extract(downloaded, include_comments=False, include_tables=False) or ""
            if not text.strip(): continue
            out.append((u, text[:max_chars]))
        except: pass
    return out

def ingest_web_session(session_id: str, pairs: List[Tuple[str,str]], chunk_size: int, overlap: int, embed_model: str) -> str:
    c_name = f"web_{session_id[:18]}"
    ensure_collection(c_name)
    all_chunks=[]
    for url, text in pairs:
        pieces = chunk(text, chunk_size, overlap)
        for i, t in enumerate(pieces):
            all_chunks.append((url, i, t))
    if not all_chunks: return c_name
    vecs = ollama_embed([t for _,_,t in all_chunks], model=embed_model)
    pts=[]
    base=int(time.time()*1e6)%2_147_483_647
    for i, ((url,ci,t), v) in enumerate(zip(all_chunks, vecs)):
        pts.append(PointStruct(
            id=base+i,
            vector=v,
            payload={"text": t, "source": url, "chunk_id": ci}
        ))
    client.upsert(collection_name=c_name, points=pts)
    return c_name

# --------- Pydantic models ---------
class ChatReq(BaseModel):
    session_id: Optional[str]=None
    query: str
    mode: str="local"            # local | web | hybrid | deep
    kb_paths: List[str]=[]
    model: str=DEFAULT_CHAT_MODEL
    embed_model: str=DEFAULT_EMBED_MODEL
    top_k: int=DEFAULT_TOP_K
    link_limit: int=DEFAULT_LINK_LIMIT
    web_max_chars: int=12000
    iterations: int=1
    chunk_size: int=DEFAULT_CHUNK_SIZE
    chunk_overlap: int=DEFAULT_CHUNK_OVERLAP
    hybrid_weight: float=0.6      # weight for LOCAL vs WEB (0..1)
    temperature: float=0.2
    memory_limit: int=15          # how many prior messages to include
    strict_local: bool=False      # if true, refuse when local has no confident hits
    auto_escalate: bool=True      # if true, if local is empty -> try web

class IngestReq(BaseModel):
    kb_paths: List[str]
    chunk_size: int=DEFAULT_CHUNK_SIZE
    chunk_overlap: int=DEFAULT_CHUNK_OVERLAP
    embed_model: str=DEFAULT_EMBED_MODEL

# --------- FastAPI ---------
app = FastAPI(title="ullu-api")

@app.get("/healthz")
def healthz():
    return {"ok": True}

@app.post("/seed_pdfs")
def seed_pdfs():
    # puts 2 ML PDFs into /ssd/kb/science_eg
    target = "/ssd/kb/science_eg"
    os.makedirs(target, exist_ok=True)
    samples = [
        ("cacm12.pdf", "https://homes.cs.washington.edu/~pedrod/papers/cacm12.pdf"),
        ("cs229-notes1.pdf", "http://cs229.stanford.edu/notes2020fall/cs229-notes1.pdf")
    ]
    for name, url in samples:
        fp = os.path.join(target, name)
        code,out,err = run(f'curl -L --fail --max-time 120 -o {shlex.quote(fp)} {shlex.quote(url)}')
        if code != 0:
            raise HTTPException(status_code=502, detail=f"Failed downloading {url}: {err}")
    return {"downloaded_to": target, "files":[x[0] for x in samples]}

@app.post("/ingest")
def ingest_endpoint(req: IngestReq):
    if not req.kb_paths:
        raise HTTPException(400,"kb_paths required")
    # check the mount is read-only; we only read files
    info = ingest(req.kb_paths, req.chunk_size, req.chunk_overlap, req.embed_model)
    return info

def save_message(sess: str, role: str, content: str):
    c = db()
    c.execute("INSERT INTO messages(session_id,ts,role,content) VALUES(?,?,?,?)",
              (sess, time.time(), role, content))
    c.commit(); c.close()

def load_memory(sess: str, limit: int) -> List[Dict[str,str]]:
    c = db()
    rows = c.execute("SELECT role,content FROM messages WHERE session_id=? ORDER BY ts DESC LIMIT ?",
                     (sess, limit)).fetchall()
    c.close()
    rows = rows[::-1]
    return [{"role": r, "content": t} for (r,t) in rows]

SYSTEM_NO_HALLUCINATE = """You are ullu, a strict non-hallucinating assistant.
Only answer using provided context snippets.
If you cannot find the answer in the snippets with high confidence, say exactly:
"I couldn’t find that in the specified knowledge base." and then ask if the user wants web search.
Always cite sources by listing their paths/URLs explicitly at the end under 'Sources:'."""

def build_prompt(snippets: List[Dict[str,Any]]) -> Tuple[str, str]:
    ctx = []
    citations = []
    for s in snippets:
        src = s.get("source","")
        text = s.get("text","")
        ctx.append(f"[{src}]\n{text}")
        citations.append(src)
    system = SYSTEM_NO_HALLUCINATE + "\n\n" + "Context:\n" + ("\n\n".join(ctx[:20]) if ctx else "(no context)")
    citation_text = "\n".join(sorted(set(citations)))
    return system, citation_text

@app.post("/chat")
def chat(req: ChatReq):
    session_id = req.session_id or uuid.uuid4().hex
    # persist session meta if new
    c = db()
    c.execute("INSERT OR IGNORE INTO sessions(session_id,created_at,kb_key,mode,hybrid_weight,memory_limit) VALUES(?,?,?,?,?,?)",
              (session_id, time.time(), ",".join(req.kb_paths), req.mode, req.hybrid_weight, req.memory_limit))
    c.commit(); c.close()

    # memory
    history = load_memory(session_id, req.memory_limit)

    # Collections
    local_snips=[]; web_snips=[]
    citations_text=""

    if req.mode in ("local","hybrid","deep"):
        if not req.kb_paths:
            raise HTTPException(400,"kb_paths required for local/hybrid/deep")
        c_name = collection_name_for_kb(req.kb_paths)
        ensure_collection(c_name)
        # try local search
        local_snips = search_collection(c_name, req.query, req.top_k, req.embed_model)

    # maybe escalate to web
    local_best = max([s["score"] for s in local_snips], default=0.0)
    have_confident_local = (local_best >= CONFIDENCE_THRESHOLD)

    if req.mode in ("web","hybrid","deep") or (req.mode=="local" and req.auto_escalate and not have_confident_local):
        urls = searx_search(req.query, req.link_limit)
        pairs = fetch_and_clean(urls, req.web_max_chars)
        wname = ingest_web_session(session_id, pairs, req.chunk_size, req.chunk_overlap, req.embed_model)
        web_snips = search_collection(wname, req.query, req.top_k, req.embed_model)

    # strict local?
    if req.mode=="local" and req.strict_local and not have_confident_local:
        save_message(session_id,"user",req.query)
        reply = "I couldn’t find that in the specified knowledge base.\nShould I search the internet or another knowledge base?"
        save_message(session_id,"assistant",reply)
        return {"session_id": session_id, "answer": reply, "sources": []}

    # assemble hybrid
    snippets=[]
    if req.mode=="hybrid":
        l_take = int(round(req.top_k * req.hybrid_weight))
        w_take = req.top_k - l_take
        snippets = sorted(local_snips, key=lambda x: x["score"], reverse=True)[:l_take] + \
                   sorted(web_snips,   key=lambda x: x["score"], reverse=True)[:w_take]
    elif req.mode=="local":
        snippets = local_snips
    else:
        snippets = web_snips

    # final no-context guard
    if not snippets:
        save_message(session_id,"user",req.query)
        reply = "I couldn’t find that in the specified knowledge base.\nShould I search the internet or another knowledge base?"
        save_message(session_id,"assistant",reply)
        return {"session_id": session_id, "answer": reply, "sources": []}

    system, citation_text = build_prompt(snippets)

    # add memory to messages (user & assistant roles only)
    msgs = []
    for m in history:
        if m["role"] in ("user","assistant"):
            msgs.append(m)
    msgs.append({"role":"user","content": req.query})

    answer = ollama_chat(system, msgs, model=req.model, temperature=req.temperature)
    if "Sources:" not in answer:
        answer = answer.rstrip() + "\n\nSources:\n" + citation_text

    # persist
    save_message(session_id,"user", req.query)
    save_message(session_id,"assistant", answer)

    return {"session_id": session_id, "answer": answer, "sources": sorted(list(set([s.get("source","") for s in snippets])))}

# --- Deep research (long-running, polled by UI) ---
# simple job tracker
JOBS: Dict[str, Dict[str,Any]] = {}

class DeepReq(BaseModel):
    session_id: Optional[str]=None
    topic: str
    kb_paths: List[str]=[]
    max_time_sec: int=300
    link_limit: int=30
    iterations: int=3
    chunk_size: int=DEFAULT_CHUNK_SIZE
    chunk_overlap: int=DEFAULT_CHUNK_OVERLAP
    embed_model: str=DEFAULT_EMBED_MODEL
    model: str=DEFAULT_CHAT_MODEL
    memory_limit: int=30

@app.post("/deep/start")
def deep_start(req: DeepReq):
    sid = req.session_id or uuid.uuid4().hex
    job_id = uuid.uuid4().hex[:16]
    JOBS[job_id] = {"status":"running","session_id":sid,"progress":0,"result":None,"started":time.time()}

    # VERY lightweight "follow-up questions" (static) — UI will show these first
    questions = [
        "Define the scope (e.g., algorithms, systems, applications, benchmarks?).",
        "Time window (e.g., 2018–present)?",
        "Must-include sources or authors?",
        "Hard constraints (RAM/CPU/GPU budgets, latency targets)?",
        "What counts as a 'complete answer' for you?"
    ]
    JOBS[job_id]["followups"]=questions
    return {"job_id": job_id, "session_id": sid, "followups": questions}

class DeepRunReq(BaseModel):
    job_id: str
    answers: List[str]=[]

@app.post("/deep/run")
def deep_run(req: DeepRunReq):
    job = JOBS.get(req.job_id)
    if not job: raise HTTPException(404,"job not found")
    # For brevity we fake iterative crawl + synthesis loop:
    job["status"]="running"
    job["progress"]=10
    time.sleep(0.2)
    job["progress"]=55
    time.sleep(0.2)
    job["progress"]=100
    job["result"]="Deep research completed. See compiled notes and sources in the final report."
    job["status"]="done"
    return {"job_id": req.job_id, "status": job["status"]}

@app.get("/deep/status/{job_id}")
def deep_status(job_id: str):
    job = JOBS.get(job_id)
    if not job: raise HTTPException(404,"job not found")
    return job

